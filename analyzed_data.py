import pandas as pd

# Load the cleaned data
file_path = r'cleaned_mtn_sprice.csv'
df = pd.read_csv(file_path)

# Clean column names and strip whitespace
df.columns = df.columns.str.strip()

# Convert the Daily Date column to datetime (auto-detects ISO format)
df['Daily Date'] = pd.to_datetime(df['Daily Date'], errors='coerce')

# Extract year
df['Year'] = df['Daily Date'].dt.year


# Step 2: Group by Year and get the max and min Year High and Year Low
yearly_data = df.groupby('Year').agg({
    'Year High (GH¢)': 'max',
    'Year Low (GH¢)': 'min'
}).reset_index()

# Step 3: Calculate the growth in price (High - Low) for each year
yearly_data['Price Growth (GH¢)'] = yearly_data['Year High (GH¢)'] - yearly_data['Year Low (GH¢)']
yearly_data['Growth Percentage (%)'] = (yearly_data['Price Growth (GH¢)'] / yearly_data['Year Low (GH¢)']) * 100

# Display the data
print(yearly_data)


import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# ----- Chart 1: Yearly High vs Low -----
plt.figure(figsize=(10, 5))
plt.bar(yearly_data['Year'] - 0.2, yearly_data['Year High (GH¢)'], width=0.4, label='Year High', color='green')
plt.bar(yearly_data['Year'] + 0.2, yearly_data['Year Low (GH¢)'], width=0.4, label='Year Low', color='red')
plt.xlabel('Year')
plt.ylabel('Price (GH¢)')
plt.title('MTN Ghana Share Price: Yearly High vs Low')
plt.legend()
plt.tight_layout()
plt.savefig("yearly_high_low.png")
plt.close()

# ----- Chart 2: Growth Percentage -----
# ----- Chart 2: Growth Percentage -----
plt.figure(figsize=(10, 5))
plt.plot(yearly_data['Year'], yearly_data['Growth Percentage (%)'], 
         color='blue', marker='o', linestyle='-', linewidth=2)
plt.xlabel('Year')
plt.ylabel('Growth (%)')
plt.title('MTN Ghana Share Price Growth by Year')
plt.grid(True, linestyle='--', alpha=0.6)
plt.tight_layout()
plt.savefig("growth_percentage.png")
plt.close()


# ----- Create PowerPoint -----
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Slide 1 - High vs Low
slide1 = prs.slides.add_slide(blank_slide_layout)
slide1.shapes.add_picture("yearly_high_low.png", Inches(1), Inches(1), width=Inches(8))

# Slide 2 - Growth %
slide2 = prs.slides.add_slide(blank_slide_layout)
slide2.shapes.add_picture("growth_percentage.png", Inches(1), Inches(1), width=Inches(8))

# Save PowerPoint
prs.save("MTN_Growth_Analysis.pptx")



with pd.ExcelWriter("MTN_Growth_Analysis.xlsx", engine='openpyxl') as writer:
    df.to_excel(writer, sheet_name="Raw Data", index=False)
    yearly_data.to_excel(writer, sheet_name="Yearly Summary", index=False)

prs.save("MTN_Growth2222_Analysis_v2.pptx")
